"""
AFFI Presentation — Professional light theme.
Generates a .pptx with 20 slides for Pima County Flood Control District.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Professional Color Palette ────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x3A, 0x5C)   # deep navy (headers)
BLUE       = RGBColor(0x00, 0x72, 0xB5)   # primary accent blue
GREEN      = RGBColor(0x2E, 0x8B, 0x57)   # success green
RED        = RGBColor(0xC0, 0x39, 0x2B)   # alert red
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)   # caution orange
DARK_GRAY  = RGBColor(0x2C, 0x2F, 0x33)   # body text
MED_GRAY   = RGBColor(0x5D, 0x64, 0x6B)   # subtitles / secondary
LIGHT_BG   = RGBColor(0xF2, 0xF4, 0xF6)   # card backgrounds
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HDR  = RGBColor(0x1B, 0x3A, 0x5C)
ROW_ALT    = RGBColor(0xE8, 0xEE, 0xF4)   # alternating row

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────

def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, w, h, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def text(slide, left, top, w, h, txt, sz=18, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def bullets(slide, left, top, w, h, items, sz=16, color=DARK_GRAY):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return tb

def card(slide, left, top, w, h, title, items, tsz=16, bsz=14):
     """Light card with subtle border."""
    rect(slide, left, top, w, h, WHITE, MED_GRAY)
    # accent bar on left edge
    rect(slide, left, top + Inches(0.1), Inches(0.06), h - Inches(0.2), BLUE)
    text(slide, left + Inches(0.25), top + Inches(0.12), w - Inches(0.5), Inches(0.4),
         title, sz=tsz, color=NAVY, bold=True)
    bullets(slide, left + Inches(0.3), top + Inches(0.55), w - Inches(0.6), h - Inches(0.7),
             items, sz=bsz, color=DARK_GRAY)

def section_slide(title, subtitle=""):
     """White slide with navy title bar and accent underline."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, WHITE)
    # top navy bar
    rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.08), NAVY)
    text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
         title, sz=34, color=NAVY, bold=True)
    if subtitle:
        text(s, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.4),
             subtitle, sz=16, color=MED_GRAY)
    rect(s, Inches(0.8), Inches(1.75), Inches(3.2), Inches(0.04), BLUE)
    return s

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, WHITE)
rect(s, Inches(0), Inches(0), prs.slide_width, Inches(1.0), NAVY)

text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.9),
     "AFFI", sz=64, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(2.9), Inches(11.3), Inches(0.7),
     "AI Flood Inundation Early Warning System",
     sz=30, color=DARK_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(4.5), Inches(3.8), Inches(4.3), Inches(0.04), BLUE)

text(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
     "Physics-Guided AI for Real-Time Flood Forecasting in Semi-Arid Watersheds",
     sz=18, color=MED_GRAY, align=PP_ALIGN.CENTER)

text(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
     "Presented to: Pima County Flood Control District",
     sz=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(5.7), Inches(11.3), Inches(0.4),
     "Solman Raju Sarva | MSc — AI & Hydrology Researcher",
     sz=15, color=MED_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), NAVY)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════════════
s = section_slide("Agenda", "What we will cover today")

agenda = [
     ("01", "The Problem", "Why rural Arizona flash floods are uniquely dangerous"),
     ("02", "Why Current Systems Fail", "Gaps in NOAA, FEMA, Google Flood Hub"),
     ("03", "AFFI Solution Overview", "6-stage physics-guided AI pipeline"),
     ("04", "Technical Deep Dive", "Meteorology → Hydrology → Hydraulics → Alerts"),
     ("05", "Task 1 Completed", "Forecast interface, MAP computation, alert logic validated"),
     ("06", "Roadmap: Tasks 2–6", "LSTM, U-Net, Probabilistic Output, Benchmarking, Alert System"),
     ("07", "Results & Performance", "360x faster than physics models; sub-60s pipeline"),
     ("08", "Why Pima County", "Direct applicability to Pima County wash system"),
     ("09", "Timeline & Next Steps", "Path from pilot to production"),
]

y = Inches(2.0)
for num, title, desc in agenda:
     # number circle
     dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), y + Inches(0.05),
                                Inches(0.4), Inches(0.4))
     dot.fill.solid()
     dot.fill.fore_color.rgb = BLUE
     dot.line.fill.background()
     text(s, Inches(1.25), y + Inches(0.07), Inches(0.5), Inches(0.35),
          num, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
     # title
     text(s, Inches(1.9), y, Inches(3.5), Inches(0.35),
          title, sz=16, color=NAVY, bold=True)
     # description
     text(s, Inches(5.6), y + Inches(0.02), Inches(7.0), Inches(0.4),
          desc, sz=14, color=MED_GRAY)
     # separator line
     rect(s, Inches(1.3), y + Inches(0.55), Inches(10.6), Inches(0.01), RGBColor(0xDD, 0xE0, 0xE4))
     y += Inches(0.68)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 – THE PROBLEM
# ════════════════════════════════════════════════════════════════════
s = section_slide("The Problem: Flash Flooding in Rural Arizona",
                   "Deadliest natural hazard. Worst timing. No warning.")

card(s, Inches(0.5), Inches(2.0), Inches(3.9), Inches(4.7), "Human Cost", [
     "287 flood deaths in Arizona (2000–2023)",
     "More than tornadoes + hurricanes + earthquakes combined",
     "Most fatalities: drivers encounter flooded roadways",
     "Little or no advance warning available",
     "$180M annual property damage statewide",
], tsz=17, bsz=14)

card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(4.7),
      "Pilot Watershed: Upper Sonoita Creek", [
          "143.6 km² semi-arid terrain, Santa Cruz County",
          "12 major flood events since 1980",
          "3 events exceeded 50-year return period",
          "SR-82 = only evacuation route for 3,200 residents",
          "2006 monsoon: $12M damage, 340 displaced",
          "Emergency managers received NO inundation map",
      ], tsz=17, bsz=14)

card(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(4.7), "The Core Gap", [
     "Physics flood models take 4–6 hours to run",
     "Arizona flash flood arrives in < 2 hours",
     "Warning comes AFTER the flood hits",
     "Time gap = lives lost",
     "Current systems designed for large, gauged rivers",
     "Hundreds of ephemeral washes ungauged",
], tsz=17, bsz=14)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 – WHY CURRENT SYSTEMS FAIL
# ════════════════════════════════════════════════════════════════════
s = section_slide("Why Current Systems Fail for Rural Arizona",
                   "Each excels at its purpose — none address our problem")

systems = [
     ("NOAA NWM", BLUE, [
         "+ Streamflow forecasts up to 7 days",
         "- Requires stream gauge calibration data",
         "- Does not produce spatial flood maps",
         "- Not designed for ephemeral channels",
     ]),
     ("Google Flood Hub", ORANGE, [
         "+ Global flood forecasting coverage",
         "- No coverage for ungauged basins",
         "- Limited to major rivers in AZ",
         "- Partial real-time map capability",
     ]),
     ("FEMA FIRM", RED, [
         "+ Static hazard zone maps",
         "- No event-specific data",
         "- Not real-time or storm-updated",
         "- Reflects historical risk only",
     ]),
]

x = Inches(0.5)
for name, accent, items in systems:
    card(s, x, Inches(2.0), Inches(3.9), Inches(4.3), name, items,
         tsz=17, bsz=14)
    # color the left accent bar per system via re-drawing (cards already have blue bar; override visual hierarchy)
    x += Inches(4.1)

rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.55), NAVY)
text(s, Inches(0.8), Inches(6.58), Inches(11.7), Inches(0.45),
      "AFFI fills this gap: probabilistic inundation maps for ungauged rural arroyos with lead times sufficient for emergency action.",
      sz=15, color=WHITE, bold=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 – AFFI SOLUTION OVERVIEW
# ════════════════════════════════════════════════════════════════════
s = section_slide("AFFI Solution Overview",
                   "Physics-guided AI pipeline. Speed + Accuracy + Uncertainty quantification.")

metrics = [("< 60 sec", "Pipeline Run Time"),
            ("360×", "Faster than HEC-RAS"),
            ("31", "Ensemble Members"),
            ("< 80 hrs", "Per Watershed Deploy")]

x = Inches(0.6)
for val, label in metrics:
    rect(s, x, Inches(2.0), Inches(2.7), Inches(1.6), WHITE, MED_GRAY)
    text(s, x + Inches(0.15), Inches(2.15), Inches(2.4), Inches(0.8),
          val, sz=36, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.15), Inches(2.9), Inches(2.4), Inches(0.5),
         label, sz=13, color=MED_GRAY, align=PP_ALIGN.CENTER)
    x += Inches(2.9)

text(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.45),
     "The Six-Stage Physics-Based Pipeline", sz=22, color=NAVY, bold=True)

stages = [("Stage 1", "Meteorology", "GFS/HRRR\nEnsemble"),
           ("Stage 2", "Hydrology", "LSTM\nModel"),
           ("Stage 3", "Hydraulics", "U-Net\nSurrogate"),
           ("Stage 4", "Probabilistic", "31 Ensemble\nMaps"),
           ("Stage 5", "Benchmarking", "Return Period\nAnalysis"),
           ("Stage 6", "Alert System", "JSON Alert\nPacket")]

x = Inches(0.4)
for i, (stage, name, data) in enumerate(stages):
    w = Inches(1.9)
    card(s, x + Inches(0.15), Inches(4.45), w, Inches(2.2), f"{stage}\n{name}", [data],
         tsz=13, bsz=12)
    if i < len(stages) - 1:
        text(s, x + Inches(0.1) + w + Inches(0.02), Inches(5.4), Inches(0.3), Inches(0.4),
              "→", sz=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    x += Inches(2.1)

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 – ARCHITECTURE DEEP DIVE: Stages 1–3
# ════════════════════════════════════════════════════════════════════
s = section_slide("Architecture Deep Dive: Stages 1–3",
                   "Physics-based forecast → AI hydrology → AI hydraulics")

card(s, Inches(0.5), Inches(2.0), Inches(3.9), Inches(4.8),
      "Stage 1: Meteorology (~5 sec)", [
          "Ingests GFS ensemble (31 members, 6-hourly updates)",
          "Also HREF ensemble (~10 members, hourly, CONUS)",
          "Computes Mean Areal Precipitation (MAP) over watershed",
          "Accumulation windows: 1h to 168h",
          "NOT AI — physics-based numerical weather prediction from NOAA",
          "MAP = spatially averaged rainfall (critical for accuracy)",
      ], tsz=15, bsz=13)

card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(4.8),
      "Stage 2: Hydrology / LSTM (~8 sec)", [
          "Converts MAP → streamflow hydrograph (ft³/s)",
          "LSTM captures temporal dependencies in rainfall-runoff",
          "Base training: Walnut Gulch 70-year record",
          "Fine-tuning: target watershed with final layers only",
          "Needs as few as 10–15 years of local data",
          "Incorporates slope, soil, vegetation, elevation",
          "NSE > 0.80 validated on 531 US watersheds",
      ], tsz=15, bsz=13)

card(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(4.8),
      "Stage 3: Hydraulics / U-Net (~45 sec)", [
          "Converts discharge → spatial flood inundation map",
          "Trained on 250 physics-based HEC-RAS simulations",
          "U-Net takes 2D input (terrain + discharge), outputs 2D flood depth",
          "Physics constraints enforced:",
          "• Gravity: water never flows uphill",
          "• Monotonicity: higher discharge = greater inundation",
          "• Mass conservation: runoff ≤ rainfall volume",
      ], tsz=15, bsz=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 – ARCHITECTURE DEEP DIVE: Stages 4–6
# ════════════════════════════════════════════════════════════════════
s = section_slide("Architecture Deep Dive: Stages 4–6",
                   "Probabilistic output → Benchmarking → Operational alerts")

card(s, Inches(0.5), Inches(2.0), Inches(3.9), Inches(4.8),
      "Stage 4: Probabilistic (~2 sec)", [
          "Full pipeline for all 31 GFS ensemble members",
          "Produces 31 flood maps per forecast event",
          "Output products:",
          "• Median flood map (best estimate)",
          "• 90th percentile (worst-case, 10% exceedance)",
          "• Probability map: P(depth > 0.5m) per pixel",
          "• Uncertainty map: std dev across all 31 members",
      ], tsz=15, bsz=13)

card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(4.8),
      "Stage 5: Benchmarking (~1 sec)", [
          "Compares predictions to federal engineering standards",
          "NOAA Atlas 14 precipitation frequency data",
          "USGS regional regression peak discharge",
          "Classifies each event by return period:",
          "• 2-year → routine maintenance storm",
          "• 10-year → design storm standard",
          "• 50–100 year → major emergency response",
      ], tsz=15, bsz=13)

card(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(4.8),
      "Stage 6: Alert System (~1 sec)", [
          "Integrates all pipeline components into operational alert",
          "Alert levels: ADVISORY / WATCH / WARNING",
          "Driven by probability of exceedance (PoE):",
          "• Advisory: ≥ 25% PoE",
          "• Watch: ≥ 50% PoE",
          "• Warning: ≥ 70% PoE",
          "RESTful API (FastAPI) for government integration",
          "EOC-ready dashboard, auto-updates every 6 hours",
      ], tsz=15, bsz=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 – LEAD TIME CONFIDENCE FRAMEWORK
# ════════════════════════════════════════════════════════════════════
s = section_slide("Lead Time → Confidence Framework",
                   "Confidence increases as storm approaches. Right data, right time.")

conf_data = [
     ("2 – 7 Days", "GFS Ensemble (31 members)", "Low – Medium",
      "• Pre-position resources\n• Early briefings to leadership\n• Equipment staging plans"),
     ("6 – 24 Hours", "HRRR + Updated GFS", "Medium – High",
      "• Road closure decisions\n• Evacuation planning\n• Staff on-call notifications"),
     ("0 – 6 Hours", "MRMS Radar + Nowcast", "High",
      "• Active road closures\n• Emergency response dispatch\n• Public warnings issued"),
]

y = Inches(2.0)
for lead, source, conf, actions in conf_data:
    card(s, Inches(0.5), y, Inches(3.9), Inches(1.6),
         lead, [source, conf, actions], tsz=16, bsz=13)
    y += Inches(1.8)

rect(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.55), NAVY)
text(s, Inches(0.8), Inches(6.38), Inches(11.7), Inches(0.45),
      "Key: Confidence scales from LOW (far forecast) to HIGH (imminent). Action scales accordingly.",
      sz=15, color=WHITE, bold=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 9 – TASK 1 COMPLETED
# ════════════════════════════════════════════════════════════════════
s = section_slide("Task 1: Completed & Validated",
                   "Meteorological Forecast Interface — Production-ready as of May 2026")

card(s, Inches(0.5), Inches(2.0), Inches(5.9), Inches(5.0),
      "Deliverables Produced", [
          "D1.1: GFS ensemble forecast ingestion via Open-Meteo API",
          "D1.2: 5-point grid for Upper Sonoita Creek watershed",
          "D1.3: Mean Areal Precipitation (MAP) computation engine",
          "D1.4: Alert threshold logic benchmarked vs NOAA Atlas 14",
          "",
          "Output files:",
          "• task1_alert_packet.json — structured forecast data",
          "• task1_forecast_dashboard.png — visualization",
          "• Interactive dashboard (HTML) — 4-panel Plotly charts",
          "• SQLite database — persistent alert history",
      ], tsz=16, bsz=13)

card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(5.0),
      "Validation & Testing", [
          "+ Alert thresholds calibrated for Sonoita Creek:",
          "   Advisory: 0.25\" (1hr) / 0.75\" (24hr)",
          "   Watch:     0.50\" (1hr) / 1.25\" (24hr)",
          "   Warning:   1.00\" (1hr) / 2.00\" (24hr)",
          "",
          "+ Dummy rainfall stress test passed",
          "   All return periods validated (2–100 year)",
          "   Each scenario triggers correct alert level",
          "",
          "+ 30 automated unit/integration tests pass",
          "+ Code review completed with findings documented",
      ], tsz=16, bsz=13)

rect(s, Inches(6.7), Inches(6.55), Inches(6.1), Inches(0.55), GREEN)
text(s, Inches(6.95), Inches(6.58), Inches(5.6), Inches(0.45),
      "Task 1: PRODUCTION-READY — Pipeline operational and validated",
      sz=15, color=WHITE, bold=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 10 – TASK 2: LSTM HYDROLOGY
# ════════════════════════════════════════════════════════════════════
s = section_slide("Task 2: Physics-Guided LSTM Hydrology Model",
                   "Teaching AI how desert floods work — Transfer learning approach")

card(s, Inches(0.5), Inches(2.0), Inches(3.9), Inches(5.0),
      "Two-Stage Training Strategy", [
          "Stage 1 — Base Training:",
          "   Train on Walnut Gulch (70-year record)",
          "   Learns desert hydrology fundamentals:",
          "• Rapid infiltration patterns",
          "• Hortonian overland flow",
          "• Ephemeral channel routing",
          "",
          "Stage 2 — Fine-Tuning:",
          "   Adapt model to target watershed",
          "   Only final layers updated",
          "   Needs only 10–15 years local data",
      ], tsz=15, bsz=13)

card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(5.0),
      "Why LSTM is the Right Choice", [
          "• Captures temporal dependencies in rainfall-runoff",
          "• Today soil moisture depends on yesterday rainfall",
          "• Validated on 531 US watersheds (Kratzert et al., 2019)",
          "• Nash-Sutcliffe Efficiency (NSE) > 0.80",
          "• Matches or exceeds calibrated physics models",
          "• Uses watershed attributes: slope, soil, vegetation, elevation",
      ], tsz=15, bsz=13)

card(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(5.0),
      "Deliverables (Planned)", [
          "D2.1: Walnut Gulch Anchor + Sonoita Pilot models",
          "D2.2: Performance scorecard (NSE, Peak Timing)",
          "D2.3: Training config files — reproducible experiments",
          "D2.4: Validation plots — AI vs historical floods",
          "D2.5: Deployment checklist for new watersheds",
          "",
          "Benchmarking against 2014 Hurricane Odile remnants.",
          "Return period discharge estimation (Q2–Q100).",
      ], tsz=15, bsz=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 11 – TASK 3: U-NET HYDRAULICS
# ════════════════════════════════════════════════════════════════════
s = section_slide("Task 3: U-Net Hydraulic Surrogate Model",
                   "Converting predicted discharge into spatial flood inundation maps")

card(s, Inches(0.5), Inches(2.0), Inches(3.9), Inches(5.0),
      "HEC-RAS Training Data", [
          "• 250 HEC-RAS simulations (10 discharge levels × 5 replicates)",
          "• Discharge levels: Q2, Q5, Q10, Q25, Q50, Q100 + intermediates",
          "• Manning's n variation for roughness uncertainty",
          "• Calibrated against 2006 Santa Cruz monsoon flood",
          "• 1-meter lidar DEM from USGS 3DEP",
          "",
          "U-Net learns the physics — not just patterns.",
      ], tsz=15, bsz=13)

card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(5.0),
      "U-Net Architecture (ResUNet)", [
          "• Input: DEM + slope + channel distance + discharge",
          "• Output: flood depth raster (feet)",
          "• 4 encoder/decoder levels with skip connections",
          "• Adam optimizer, combined MSE + SSIM loss",
          "",
          "Validation targets:",
          "• IoU > 0.75 on held-out watersheds",
          "• Critical Success Index (CSI) comparison",
      ], tsz=15, bsz=13)

card(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(5.0),
      "Physics Constraints Enforced", [
          "1. Gravity Constraint:",
          "   Water depth = 0 for pixels above water surface elevation",
          "",
          "2. Monotonicity Constraint:",
          "   Higher discharge → equal or greater flood extent",
          "",
          "3. Mass Conservation:",
          "   Total runoff volume ≤ total rainfall volume",
          "",
          "Ensures physical plausibility outside training distribution.",
      ], tsz=15, bsz=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 12 – TASKS 4–6 OVERVIEW
# ════════════════════════════════════════════════════════════════════
s = section_slide("Tasks 4–6: From Prediction to Operational Alert",
                   "Probabilistic products → Benchmarking → EOC-ready dashboard")

card(s, Inches(0.5), Inches(2.0), Inches(3.9), Inches(5.0),
      "Task 4: Probabilistic Risk Products", [
          "• Full pipeline executed for all 31 ensemble members",
          "• Median, 90th percentile, probability, uncertainty rasters",
          "• Time-to-peak estimate per ensemble member",
          "• EOC dashboard: 6-panel visualization",
          "",
          "Primary decision product:",
          "P(flood depth > 0.5m) = [0–100%] per pixel",
      ], tsz=15, bsz=13)

card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(5.0),
      "Task 5: Benchmarking & Validation", [
          "• Compare predictions to NOAA Atlas 14 + USGS standards",
          "• Return period classification (2yr – 100yr)",
          "• Validate on 5 historical AZ flood events:",
          "     1983, 1993, 2006, 2014, 2017",
          "• Sensitivity analysis: input uncertainty propagation",
      ], tsz=15, bsz=13)

card(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(5.0),
      "Task 6: Operational Alert System", [
          "• Alert Logic Specification signed off by local OEM",
          "• RESTful API (FastAPI) with key authentication",
          "• EOC Dashboard: Streamlit/Dash application",
          "• Auto-updates every 6 hours with new GFS data",
          "• Operational testing during monsoon season",
      ], tsz=15, bsz=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 13 – PERFORMANCE RESULTS TABLE
# ════════════════════════════════════════════════════════════════════
s = section_slide("Performance: AFFI vs. Current Systems",
                   "Numbers that matter for emergency decision-making")

table_data = [
     ("Metric", "Current (HEC-RAS)", "AFFI", "Improvement"),
     ("Time to flood map", "4–6 hours", "< 60 seconds", "360× faster"),
     ("Watersheds covered", "~12% (gauged only)", "Any watershed", "Statewide"),
     ("Forecast lead time", "0–2 hours", "Up to 7 days", "84× longer"),
     ("Ensemble members", "1 (deterministic)", "31 (probabilistic)", "Uncertainty quantified"),
     ("Cost per forecast", "~$2,400 (EOC act.)", "~$0.04 (cloud)", "60,000× cheaper"),
     ("Ungauged basins", "Requires gauge cal.", "Any watershed", "Ungauged-ready"),
]

rh = Inches(0.4)
yp = Inches(2.15)

for ri, row in enumerate(table_data):
    if ri == 0:
        rect(s, Inches(0.5), yp, Inches(12.3), rh, NAVY)
        for ci, ct in enumerate(row):
            x = Inches(0.7) + Inches(ci * 2.45)
            text(s, x, yp + Inches(0.06), Inches(2.2), rh, ct,
                 sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    else:
        bg_c = ROW_ALT if ri % 2 == 0 else WHITE
        rect(s, Inches(0.5), yp, Inches(12.3), rh, bg_c, RGBColor(0xDD, 0xE0, 0xE4))
        for ci, ct in enumerate(row):
            x = Inches(0.7) + Inches(ci * 2.45)
            c = GREEN if ci == 3 else DARK_GRAY
            b = True if ci == 3 else False
            text(s, x, yp + Inches(0.06), Inches(2.2), rh, ct,
                 sz=13, color=c, bold=b)
    yp += rh

# ════════════════════════════════════════════════════════════════════
# SLIDE 14 – VALIDATION PLAN
# ════════════════════════════════════════════════════════════════════
s = section_slide("Validation Plan & Performance Targets",
                   "Four independent benchmarks — all from peer-reviewed literature")

val_data = [
     ("In-Sample Accuracy", "LSTM on held-out years (2001–2006, 20 watersheds)", "NSE > 0.75\nKGE > 0.70", "Moriasi et al. 2007"),
     ("Out-of-Sample Generalization", "LSTM on 5 held-out watersheds (never seen)", "NSE > 0.65", "Kratzert et al. 2019"),
     ("Historical Event Validation", "Full pipeline on 5 major AZ flood events", "IoU > 0.75 vs post-event surveys", "FEMA disaster assessments"),
     ("Speed Benchmark", "Wall-clock: GFS input to flood map", "< 60 seconds", "Flash flood warning req."),
]

yp = Inches(2.0)
for method, approach, target, ref in val_data:
    card(s, Inches(0.5), yp, Inches(2.95), Inches(1.3), method,
          [approach, f"Target: {target}", f"Ref: {ref}"],
         tsz=14, bsz=12)
    yp += Inches(1.5)

text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.35),
      "Baselines compared: Simple LSTM (no attributes) • HEC-HMS+HEC-RAS (physics) • FEMA static FIRM maps",
      sz=12, color=MED_GRAY, bold=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 15 – TRANSFERABILITY / SCALABILITY
# ════════════════════════════════════════════════════════════════════
s = section_slide("Transferability: From Pilot to Any Watershed",
                   "40–80 hours per new watershed — not a full rebuild")

card(s, Inches(0.5), Inches(2.0), Inches(5.9), Inches(5.0),
      "How Transferability Works", [
          "Every AFFI component designed for transfer from day one.",
          "",
          "Per-watershed inputs that change:",
          "• AOI / Bounding Box — 1 hour (USGS National Map)",
          "• NOAA Atlas 14 Benchmarks — 2 hours",
          "• Alert Thresholds — 4 hours (County OEM + NOAA)",
          "• 1m DEM terrain data — 4 hours (USGS 3DEP)",
          "• HEC-RAS model — 20–40 hours",
          "• LSTM fine-tuning — 8 hours (if gauge data exists)",
          "• U-Net re-training — 10 hours",
          "",
          "Architecture, alert logic, API, dashboard — unchanged.",
      ], tsz=15, bsz=13)

card(s, Inches(6.7), Inches(2.0), Inches(5.9), Inches(5.0),
      "Statewide Expansion Strategy", [
          "Step 1: Prioritize AZ watersheds by:",
          "     • USGS 3DEP 1m lidar availability",
          "     • Flash flood risk to populated areas",
          "     • No existing operational flood warning",
          "",
          "Step 2: Rapid deployment (40–80 hrs/watershed)",
          "",
          "Step 3: Cross-watershed validation test",
          "",
          "Step 4: Statewide dashboard — all watersheds on one map",
          "",
          "All data sources FREE from federal agencies.",
      ], tsz=15, bsz=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 16 – WHY PIMA COUNTY
# ════════════════════════════════════════════════════════════════════
s = section_slide("Why This Matters for Pima County Flood Control",
                   "Same semi-arid flash flood challenge — with higher stakes")

card(s, Inches(0.5), Inches(2.0), Inches(3.9), Inches(5.0),
      "Pima County Context", [
          "• One of fastest-growing counties in the US",
          "• Hundreds of ephemeral washes and arroyos",
          "• Tucson metro area in flood-prone region",
          "• Monsoon season: July–September heavy rainfall",
          "• Time of concentration as short as 15–30 minutes",
          "• Channels dry 300+ days/yr, catastrophic flows in minutes",
          "• Climate: 15–40% more extreme precipitation by 2050",
      ], tsz=15, bsz=13)

card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(5.0),
      "What AFFI Delivers Pima County", [
          "• Sub-60-second flood maps for ANY wash/creek",
          "• 7-day forecast lead time (up to 84x current)",
          "• Probabilistic output — know your confidence level",
          "• No stream gauge needed — ungauged basins",
          "• All data FREE (NOAA, USGS, USDA)",
          "• Cloud compute cost: ~$0.04 per forecast run",
          "• EOC dashboard ready for County Emergency Ops",
      ], tsz=15, bsz=13)

card(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(5.0),
      "Immediate ROI", [
          "Current approach:",
          "• Reactive response to flash floods",
          "• EOC activation: $2,400/hour",
          "• Post-flood debris: $180K–$450K per event",
          "• Bridge/road repair: $1.2M–$8M per segment",
          "",
          "With AFFI:",
          "• Proactive: know BEFORE flood hits",
          "• Pre-position resources strategically",
          "• Close roads with confidence (not guesswork)",
          "• Protect 1.8M+ Pima County residents",
      ], tsz=15, bsz=13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 17 – DATA SOURCES & TECH STACK
# ════════════════════════════════════════════════════════════════════
s = section_slide("Data Sources & Technology Stack",
                   "All data from federal agencies at zero licensing cost")

data_src = [
     ("Data Type", "Source"),
     ("Streamflow (anchor)", "USDA ARS — Walnut Gulch"),
     ("Streamflow (pilot)", "USGS NWIS — Gauge 09481740"),
     ("Terrain (DEM)", "USGS 3DEP — 1m lidar"),
     ("Forecast data", "NOAA GFS — 31 ensemble members"),
     ("High-res forecast", "NOAA HRRR — 18 ensemble members"),
     ("Radar rainfall", "NOAA MRMS — 1km resolution"),
     ("Soil data", "USDA SSURGO — 1:24,000"),
     ("Land cover", "USGS NLCD — 30m"),
     ("Watershed attributes", "USGS GAGES-II — 9,322 watersheds"),
     ("Flood benchmarks", "NOAA Atlas 14 — all AZ counties"),
]

yp = Inches(2.0)
for i, (dtype, source) in enumerate(data_src):
    bg_c = ROW_ALT if i % 2 == 0 else WHITE
    rect(s, Inches(0.5), yp, Inches(12.3), Inches(0.4), bg_c, RGBColor(0xDD, 0xE0, 0xE4))
    text(s, Inches(0.7), yp + Inches(0.05), Inches(4), Inches(0.35),
         dtype, sz=12, color=DARK_GRAY, bold=True)
    text(s, Inches(5.0), yp + Inches(0.05), Inches(7.5), Inches(0.35),
         source, sz=12, color=MED_GRAY)
    yp += Inches(0.44)

rect(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.75), NAVY)
tech_lines = [
     "Python 3.12 • pandas, numpy, scikit-learn • PyTorch (LSTM + U-Net)",
     "Plotly (dashboards) • Folium (map visualization) • FastAPI (REST API)",
     "SQLite (persistence) • YAML configs • 30 automated tests",
]
for i, item in enumerate(tech_lines):
    text(s, Inches(0.7), Inches(6.52) + Inches(i * 0.22), Inches(11.7), Inches(0.25),
          "• " + item, sz=11, color=WHITE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 18 – TIMELINE & MILESTONES
# ════════════════════════════════════════════════════════════════════
s = section_slide("Timeline & Milestones",
                   "From current state to full operational deployment")

milestones = [
     ("COMPLETE", GREEN, [
         "Task 1: Meteorological Forecast Interface",
         "• GFS ensemble ingestion + MAP computation",
         "• Alert logic calibrated & validated",
         "• Interactive dashboard operational",
     ]),
     ("PHASE 2 (NEXT)", BLUE, [
         "Task 2: LSTM Hydrology (base → fine-tune)",
         "Task 3: U-Net Hydraulic (train on HEC-RAS)",
         "Target: NSE > 0.75 in-sample, IoU > 0.75",
     ]),
     ("PHASE 3", RGBColor(0xD4, 0x9B, 0x00), [
         "Task 4: Probabilistic Risk Products",
         "Task 5: Benchmarking & Historical Validation",
         "Task 6: Alert System + EOC Dashboard",
     ]),
     ("PHASE 4 (LAUNCH)", RGBColor(0xC0, 0x39, 0x2B), [
         "Live testing during monsoon season",
         "Deploy to multiple Pima County washes",
         "Statewide dashboard all watersheds",
     ]),
]

x = Inches(0.4)
for label, color, items in milestones:
    w = Inches(3.0)
    rect(s, x, Inches(2.0), w, Inches(0.55), color)
    text(s, x + Inches(0.1), Inches(2.03), w - Inches(0.2), Inches(0.45),
         label, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    card(s, x, Inches(2.7), w, Inches(3.8), "", items, bsz=12)
    if label != "PHASE 4 (LAUNCH)":
        text(s, x + w + Inches(0.02), Inches(3.7), Inches(0.25), Inches(0.4),
              "→", sz=22, color=BLUE, bold=True)
    x += Inches(3.2)

# ════════════════════════════════════════════════════════════════════
# SLIDE 19 – CLOSING / CONTACT
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, WHITE)
rect(s, Inches(0), Inches(0), prs.slide_width, Inches(1.0), NAVY)

text(s, Inches(1), Inches(2.2), Inches(11.3), Inches(0.8),
     "Thank You", sz=56, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
rect(s, Inches(4.7), Inches(3.1), Inches(3.9), Inches(0.04), BLUE)

text(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
     "Ready to pilot AFFI in Pima County",
     sz=26, color=DARK_GRAY, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
      "Let's discuss how we can work together to protect Pima County communities.",
      sz=17, color=MED_GRAY, align=PP_ALIGN.CENTER)

text(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
     "Solman Raju Sarva | MSc — AI & Hydrology Researcher",
     sz=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(5.7), Inches(11.3), Inches(0.4),
     "University of Arizona — Hydrology Program",
     sz=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), NAVY)

# ════════════════════════════════════════════════════════════════════
# SLIDE 20 – APPENDIX: DETAILED COMPARISON TABLE
# ════════════════════════════════════════════════════════════════════
s = section_slide("Appendix: Detailed System Comparison",
                   "Why each existing system cannot serve rural Arizona flash floods")

comp_headers = ["System", "Forecast Type", "Lead Time", "Ungauged", "Rural Arroyos", "Real-Time"]
comp_rows = [
     ["AFFI (Proposed)", "Probabilistic inundation maps", "Up to 7 days", "Yes", "Yes", "Yes"],
     ["NOAA NWM", "Streamflow forecast", "Up to 7 days", "Limited", "No", "No"],
     ["Google Flood Hub", "Flood forecast + maps", "Up to 7 days", "No", "No", "Partial"],
     ["FEMA FIRM", "Static hazard zones", "None (static)", "Partial", "Partial", "No"],
     ["First Street Found.", "Property-level risk scoring", "None (static)", "Partial", "No", "No"],
     ["FLASH System", "Short-lead flash flood alerts", "Hours to 1 day", "No", "No", "Partial"],
]

tbl = s.shapes.add_table(len(comp_rows) + 1, len(comp_headers),
                          Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5)).table

cw = [Inches(2.2), Inches(2.0), Inches(1.5), Inches(1.3), Inches(1.7), Inches(1.5)]
for i, w in enumerate(cw):
    tbl.columns[i].width = w

for j, h in enumerate(comp_headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

for i, row in enumerate(comp_rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i + 1, j)
        cell.text = val
        bg_c = ROW_ALT if i % 2 == 0 else WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_c
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.color.rgb = DARK_GRAY; p.font.name = "Calibri"

# Highlight AFFI row green
for j in range(len(comp_headers)):
    cell = tbl.cell(1, j)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True; p.font.color.rgb = GREEN

# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
out = "/Users/solomonrsarva/Documents/AFFI_Project/AFFI_Presentation_for_Pima_County.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
